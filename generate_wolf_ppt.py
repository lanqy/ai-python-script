#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
狼来了故事PPT生成脚本
专为小朋友演讲设计
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_wolf_ppt():
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 定义颜色方案
    colors = {
        'primary': RGBColor(70, 130, 180),    # 钢蓝色
        'accent': RGBColor(255, 140, 0),      # 深橙色
        'title': RGBColor(255, 69, 0),        # 红橙色
        'text': RGBColor(0, 0, 0),            # 黑色
        'background': RGBColor(255, 248, 220) # 米色
    }

    # 创建标题幻灯片
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景形状
    bg_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    # 标题
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "🐺 狼来了 🐺"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['title']
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "一个关于诚实的重要故事"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = colors['primary']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 演讲者信息
    info_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    info_frame = info_box.text_frame
    info_frame.text = "演讲人：__________"
    info_para = info_frame.paragraphs[0]
    info_para.font.size = Pt(24)
    info_para.alignment = PP_ALIGN.CENTER

    # 第2页：故事背景
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景形状
    bg_shape = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    # 标题栏
    header_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "📖 故事开始"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = colors['title']
    header_para.alignment = PP_ALIGN.CENTER
    
    # 内容
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "很久很久以前，有一个小男孩。",
        "他负责在山坡上放羊。",
        "每天，他都要看管着羊群。",
        "小男孩觉得放羊很无聊...",
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(32)
        p.space_after = Pt(20)

    # 第3页：第一次撒谎
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    header_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "第一次撒谎"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = colors['accent']
    header_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "有一天，小男孩想：",
        '"哈哈！让我开个玩笑吧！"',
        "于是他大喊：",
        '"狼来了！狼来了！救命啊！"',
        "",
        "村民们听到呼救声，",
        "急忙拿着工具跑上山。",
        "结果发现根本没有狼！",
        "小男孩哈哈大笑..."
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        if "大喊" in line or "狼来了" in line:
            p.font.color.rgb = colors['title']
            p.font.bold = True
        p.space_after = Pt(15)

    # 第4页：第二次撒谎
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    header_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "第二次撒谎"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = colors['accent']
    header_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "过了几天，小男孩又觉得无聊了。",
        "他又一次大喊：",
        '"狼来了！狼来了！"',
        "",
        "村民们又急忙跑上山。",
        "结果又被骗了！",
        "",
        "村民们很生气地说：",
        '"下次我们不会再来了！"'
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        if "狼来了" in line or "很生气" in line:
            p.font.color.rgb = colors['title']
            p.font.bold = True
        p.space_after = Pt(15)

    # 第5页：狼真的来了
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(255, 228, 225)  # 浅红色背景
    bg_shape.line.fill.background()
    
    header_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "😱 狼真的来了！"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = RGBColor(139, 0, 0)  # 深红色
    header_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "有一天，狼真的来了！",
        "🐺 一只大灰狼出现在山坡上！",
        "",
        "小男孩吓坏了，拼命大喊：",
        '"狼来了！狼来了！救命啊！"',
        "",
        "村民们听到了呼救声。",
        "但是，他们以为小男孩又在撒谎...",
        "没有人上山帮助他..."
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        p.space_after = Pt(15)

    # 第6页：后果
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    header_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "😢 后果"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = colors['title']
    header_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "小男孩很伤心，也很后悔。",
        "他哭着说：",
        '"对不起！我再也不撒谎了！"',
        "",
        "从那以后，小男孩学会了",
        "做一个诚实的人。",
        "",
        "虽然他吸取了教训，",
        "但已经太晚了..."
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(28)
        p.space_after = Pt(15)

    # 第7页：故事寓意
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(224, 255, 255)  # 浅青色背景
    bg_shape.line.fill.background()
    
    header_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    header_frame = header_box.text_frame
    header_frame.text = "💡 故事告诉我们"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = colors['primary']
    header_para.alignment = PP_ALIGN.CENTER
    
    content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    lines = [
        "1. 要做一个诚实的人",
        "",
        "2. 撒谎会失去别人的信任",
        "",
        "3. 信任一旦失去，很难再找回",
        "",
        "4. 不要拿别人的关心开玩笑",
        "",
        "5. 诚实是珍贵的品质"
    ]
    
    for line in lines:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(32)
        p.font.color.rgb = colors['text']
        p.space_after = Pt(20)

    # 第8页：结束页
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors['background']
    bg_shape.line.fill.background()
    
    title_box = slide8.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "谢谢大家！"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['title']
    title_para.alignment = PP_ALIGN.CENTER
    
    # 寓语总结
    lesson_box = slide8.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.5))
    lesson_frame = lesson_box.text_frame
    lesson_frame.text = '"诚实是一种美德，让我们每个人都做诚实的孩子吧！"'
    lesson_para = lesson_frame.paragraphs[0]
    lesson_para.font.size = Pt(28)
    lesson_para.font.color.rgb = colors['primary']
    lesson_para.alignment = PP_ALIGN.CENTER
    
    # 装饰图形
    star_box = slide8.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    star_frame = star_box.text_frame
    star_frame.text = "⭐ ⭐ ⭐ ⭐ ⭐"
    star_para = star_frame.paragraphs[0]
    star_para.font.size = Pt(36)
    star_para.alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    print("正在生成PPT...")
    prs = create_wolf_ppt()
    
    # 保存PPT
    output_path = "c:/Users/220106006/WorkBuddy/Claw/狼来了故事_演讲.pptx"
    prs.save(output_path)
    
    print("PPT已成功生成！")
    print("文件位置：" + output_path)
    print("共 " + str(len(prs.slides)) + " 页")
